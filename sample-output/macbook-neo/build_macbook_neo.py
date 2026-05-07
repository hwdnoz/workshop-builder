"""
build_macbook_neo.py
Generates macbook_neo_presentation.pptx — a 5-slide workshop deck on the MacBook Neo.
Run: python3 build_macbook_neo.py

5-slide structure:
  1. Title slide
  2. What is the MacBook Neo? (overview / key specs)
  3. Key Features (design, performance, standout capabilities)
  4. Who is it for? (use cases, audience)
  5. Key Takeaways / Closing
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---------------------------------------------------------------------------
# Color palette (3-5 hex colors, consistent throughout)
# ---------------------------------------------------------------------------
C_BG_DARK   = RGBColor(0x1A, 0x1A, 0x2E)   # deep navy — primary background
C_BG_MID    = RGBColor(0x22, 0x28, 0x3A)   # mid navy — alternate background
C_ACCENT    = RGBColor(0x00, 0x7A, 0xFF)   # Apple-style blue — headings & accents
C_BODY      = RGBColor(0xD8, 0xE0, 0xEC)   # off-white — body text
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)   # pure white — title text
C_IMG_BG    = RGBColor(0x0D, 0x1E, 0x36)   # dark panel — image placeholder fill

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])  # fully blank


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width_pt=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        if line_width_pt:
            shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, left, top, width, height,
                font_size=18, bold=False, italic=False,
                color=C_BODY, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullets(slide, bullets, left, top, width, height, font_size=18, color=C_BODY):
    """Add a bulleted list (max 4 items) as separate paragraphs in one textbox."""
    assert len(bullets) <= 4, f"Max 4 bullets per slide, got {len(bullets)}"
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"•  {bullet}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = False
    return txBox


def add_slide_number(slide, num, total=5):
    add_textbox(slide, f"{num} / {total}",
                left=Inches(12.0), top=Inches(7.05),
                width=Inches(1.1), height=Inches(0.35),
                font_size=11, color=C_ACCENT, align=PP_ALIGN.RIGHT)


def add_section_label(slide, text):
    """Movement label — top left of content slides."""
    add_textbox(slide, text.upper(),
                left=Inches(0.38), top=Inches(0.18),
                width=Inches(4.5), height=Inches(0.3),
                font_size=10, bold=True, color=C_ACCENT)


def image_placeholder(slide, left, top, width, height, label):
    """Large image placeholder — at least half the slide width."""
    add_rect(slide, left, top, width, height,
             fill_color=C_IMG_BG, line_color=C_ACCENT, line_width_pt=1.5)
    add_textbox(slide, f"[IMAGE: {label}]",
                left=left + Inches(0.2),
                top=top + (height // 2) - Pt(14),
                width=width - Inches(0.4),
                height=Inches(0.55),
                font_size=13, italic=True, color=C_ACCENT,
                align=PP_ALIGN.CENTER)


def bg_fill(slide, color=C_BG_DARK):
    add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, fill_color=color)


def accent_rule(slide, left, top, width_in=3.5):
    add_rect(slide, left, top, Inches(width_in), Pt(2.5), fill_color=C_ACCENT)


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ---------------------------------------------------------------------------
# Slide 1 — Title
# ---------------------------------------------------------------------------
def build_slide_01(prs):
    slide = blank_slide(prs)
    bg_fill(slide, C_BG_DARK)

    # Thin accent bar on left edge
    add_rect(slide, Inches(0), Inches(0), Inches(0.07), SLIDE_H, fill_color=C_ACCENT)

    # Right-side image placeholder (large — covers ~52% of slide width)
    image_placeholder(slide,
                      left=Inches(6.5), top=Inches(0.45),
                      width=Inches(6.6), height=Inches(6.6),
                      label="MacBook Neo — product hero shot on a clean surface")

    # Title
    add_textbox(slide, "MacBook Neo",
                left=Inches(0.5), top=Inches(1.8),
                width=Inches(5.8), height=Inches(1.3),
                font_size=54, bold=True, color=C_WHITE,
                align=PP_ALIGN.LEFT)

    # Descriptor line
    add_textbox(slide, "A concept at the edge of Apple's MacBook lineup",
                left=Inches(0.5), top=Inches(3.2),
                width=Inches(5.8), height=Inches(0.6),
                font_size=22, color=C_ACCENT, align=PP_ALIGN.LEFT)

    # Date
    add_textbox(slide, "May 2026",
                left=Inches(0.5), top=Inches(6.8),
                width=Inches(3.0), height=Inches(0.4),
                font_size=13, color=C_BODY, align=PP_ALIGN.LEFT)

    add_slide_number(slide, 1)

    add_notes(slide,
        "Welcome. Today we're looking at the MacBook Neo — a conceptual device that "
        "synthesizes where Apple's laptop line is heading. The 'Neo' name isn't official; "
        "think of it as a way to explore real trends in Apple Silicon, design, and on-device "
        "AI all at once. We'll cover what it is, what makes it interesting, who it's built "
        "for, and what it means for the broader MacBook lineup.")


# ---------------------------------------------------------------------------
# Slide 2 — What is the MacBook Neo? (overview / key specs)
# ---------------------------------------------------------------------------
def build_slide_02(prs):
    slide = blank_slide(prs)
    bg_fill(slide, C_BG_MID)
    add_section_label(slide, "Overview")

    add_textbox(slide, "What is the MacBook Neo?",
                left=Inches(0.4), top=Inches(0.62),
                width=Inches(7.5), height=Inches(0.85),
                font_size=34, bold=True, color=C_WHITE)

    accent_rule(slide, Inches(0.4), Inches(1.53))

    bullets = [
        "Conceptual: Air portability + Pro performance",
        "M4 Pro chip — up to 24-core GPU",
        "14\" and 16\" — enters at Pro price territory",
        "OLED display — a first for the MacBook line",
    ]
    add_bullets(slide, bullets,
                left=Inches(0.4), top=Inches(1.75),
                width=Inches(6.4), height=Inches(3.6),
                font_size=20)

    # Image placeholder — right side (~50% of slide width)
    image_placeholder(slide,
                      left=Inches(6.6), top=Inches(0.75),
                      width=Inches(6.5), height=Inches(5.9),
                      label="MacBook Neo — 14\" and 16\" models side by side")

    add_slide_number(slide, 2)

    add_notes(slide,
        "So the Neo concept merges the best of both existing lines. The MacBook Air is the "
        "thin-and-light choice — fanless, up to 24 GB, starts at $1,099. The MacBook Pro "
        "gives you M4 Pro, sustained performance, and a 120 Hz XDR display — but it's heavier. "
        "The Neo asks: what if you got the Pro chip in an Air-class chassis, added OLED, and "
        "pushed connectivity to Thunderbolt 5? That's the concept we're working through today.")


# ---------------------------------------------------------------------------
# Slide 3 — Key Features
# ---------------------------------------------------------------------------
def build_slide_03(prs):
    slide = blank_slide(prs)
    bg_fill(slide, C_BG_DARK)
    add_section_label(slide, "Key Features")

    add_textbox(slide, "Key Features",
                left=Inches(0.4), top=Inches(0.62),
                width=Inches(6.0), height=Inches(0.85),
                font_size=34, bold=True, color=C_WHITE)

    accent_rule(slide, Inches(0.4), Inches(1.53))

    # Left column — Design & Display
    add_textbox(slide, "Design & Display",
                left=Inches(0.4), top=Inches(1.72),
                width=Inches(2.95), height=Inches(0.42),
                font_size=14, bold=True, color=C_ACCENT)
    left_bullets = [
        "Under 11 mm — thinnest Mac ever",
        "OLED, 120 Hz ProMotion",
        "Nano-texture glass option",
        "Fanless thermal design",
    ]
    add_bullets(slide, left_bullets,
                left=Inches(0.4), top=Inches(2.18),
                width=Inches(2.95), height=Inches(2.8),
                font_size=17)

    # Right column — Performance & Connectivity
    add_textbox(slide, "Performance & Connectivity",
                left=Inches(3.55), top=Inches(1.72),
                width=Inches(3.0), height=Inches(0.42),
                font_size=14, bold=True, color=C_ACCENT)
    right_bullets = [
        "M4 Pro — up to 24-core GPU",
        "Up to 64 GB unified memory",
        "Thunderbolt 5 + MagSafe",
        "Wi-Fi 7 and Bluetooth 5.4",
    ]
    add_bullets(slide, right_bullets,
                left=Inches(3.55), top=Inches(2.18),
                width=Inches(3.0), height=Inches(2.8),
                font_size=17)

    # Image placeholder — right side (~50% of slide width)
    image_placeholder(slide,
                      left=Inches(6.6), top=Inches(0.75),
                      width=Inches(6.5), height=Inches(5.9),
                      label="MacBook Neo open — keyboard and display detail")

    add_slide_number(slide, 3)

    add_notes(slide,
        "The thickness is the interesting bet here — fanless at M4 Pro performance is a big "
        "thermal challenge. OLED is something Apple has been rumored to bring to MacBooks for "
        "years; the concept slots it in here. Each M chip generation delivers roughly 20-25% "
        "CPU gains — so an M4 Pro in an Air-class chassis is a genuinely different device "
        "than an M3 Air. Thunderbolt 5 matters for fast external storage and docking.")


# ---------------------------------------------------------------------------
# Slide 4 — Who is it for?
# ---------------------------------------------------------------------------
def build_slide_04(prs):
    slide = blank_slide(prs)
    bg_fill(slide, C_BG_MID)
    add_section_label(slide, "Use Cases")

    add_textbox(slide, "Who is it for?",
                left=Inches(0.4), top=Inches(0.62),
                width=Inches(7.0), height=Inches(0.85),
                font_size=34, bold=True, color=C_WHITE)

    accent_rule(slide, Inches(0.4), Inches(1.53), width_in=3.0)

    bullets = [
        "Creatives who travel and need Pro GPU",
        "Developers: Xcode, local AI models, Docker",
        "Video editors wanting on-the-go performance",
        "Intel-era MacBook owners ready to upgrade",
    ]
    add_bullets(slide, bullets,
                left=Inches(0.4), top=Inches(1.75),
                width=Inches(6.5), height=Inches(3.4),
                font_size=20)

    image_placeholder(slide,
                      left=Inches(6.6), top=Inches(0.75),
                      width=Inches(6.5), height=Inches(5.9),
                      label="Person using MacBook Neo in a cafe or airport lounge")

    add_slide_number(slide, 4)

    add_notes(slide,
        "The Neo concept targets people stuck between two good-but-not-quite-right options. "
        "Creatives who travel know the pain: Air is light but struggles with heavy GPU work, "
        "Pro is capable but heavy. Developers running local AI models — Apple Intelligence "
        "runs entirely on-device on M-series chips, which is a real privacy and latency "
        "advantage. And there's still a large Intel Mac install base that hasn't made the "
        "Apple Silicon jump yet. A device like this would be a strong reason to upgrade.")


# ---------------------------------------------------------------------------
# Slide 5 — Key Takeaways / Closing
# ---------------------------------------------------------------------------
def build_slide_05(prs):
    slide = blank_slide(prs)
    bg_fill(slide, C_BG_DARK)
    add_section_label(slide, "Takeaways")

    add_textbox(slide, "Key Takeaways",
                left=Inches(0.4), top=Inches(0.62),
                width=Inches(8.0), height=Inches(0.85),
                font_size=34, bold=True, color=C_WHITE)

    accent_rule(slide, Inches(0.4), Inches(1.53))

    bullets = [
        "Thinnest Mac ever — M4 Pro inside",
        "OLED arrives on MacBook for the first time",
        "Thunderbolt 5 and Wi-Fi 7 extend its lifespan",
        "Strong upgrade path from Intel-era machines",
    ]
    add_bullets(slide, bullets,
                left=Inches(0.4), top=Inches(1.75),
                width=Inches(7.5), height=Inches(3.5),
                font_size=20)

    add_textbox(slide, "Questions?",
                left=Inches(0.4), top=Inches(5.7),
                width=Inches(4.0), height=Inches(0.65),
                font_size=26, bold=True, color=C_ACCENT)

    image_placeholder(slide,
                      left=Inches(6.6), top=Inches(0.75),
                      width=Inches(6.5), height=Inches(5.9),
                      label="MacBook Neo closed — unibody aluminum chassis detail")

    add_slide_number(slide, 5)

    add_notes(slide,
        "The MacBook Neo is a thought experiment grounded in real trends: Apple Silicon "
        "compounding at 20-25% per generation, OLED coming to MacBooks, Thunderbolt 5 "
        "landing in the lineup, Apple Intelligence running fully on-device. Whether or not "
        "Apple ships something called 'Neo,' these are the forces shaping where the MacBook "
        "line is going. Thanks everyone — happy to take questions.")


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    prs = new_prs()

    build_slide_01(prs)
    build_slide_02(prs)
    build_slide_03(prs)
    build_slide_04(prs)
    build_slide_05(prs)

    out_path = os.path.join(
        os.path.dirname(os.path.abspath(__file__)),
        "macbook_neo_presentation.pptx"
    )
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slide count: {len(prs.slides)}")


if __name__ == "__main__":
    main()
